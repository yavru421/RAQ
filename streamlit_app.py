import streamlit as st
import os
from groq import Groq
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import io
import mimetypes
from typing import Optional, Callable, Dict
import logging
import time
import pandas as pd
from PIL import Image

st.set_page_config(page_title="RAQ - Pure Groq", page_icon="🤖", layout="centered", initial_sidebar_state="collapsed")

# Add mobile-friendly optimizations
st.markdown(
    """<style>
    /* Ensure the app is responsive on mobile */
    @media (max-width: 768px) {
        .css-1d391kg { /* Streamlit main content area */
            padding: 1rem;
        }
        .css-1v3fvcr { /* Sidebar */
            display: none;
        }
    }
    </style>""",
    unsafe_allow_html=True
)

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Initialize session state
if 'document_text' not in st.session_state:
    st.session_state.document_text = ""
if 'groq_client' not in st.session_state:
    st.session_state.groq_client = None
if 'feedback' not in st.session_state:
    st.session_state.feedback = []
if 'history' not in st.session_state:
    st.session_state.history = []

# Correct the extract_text_from_pdf function to use the correct text extraction method
def extract_text_from_pdf(file):
    """Extract text from PDF file using PyPDF2."""
    try:
        file_content = file.read()
        pdf_reader = PdfReader(io.BytesIO(file_content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        st.error(f"Error reading PDF: {str(e)}")
        return ""

def extract_text_from_docx(file):
    """Extract text from DOCX file"""
    try:
        doc = Document(file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        st.error(f"Error reading DOCX: {str(e)}")
        return ""

# Updated extract_text_from_pptx to use getattr for text attribute
def extract_text_from_pptx(file: io.BytesIO) -> str:
    """Extract text from PPTX file"""
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                text += getattr(shape, "text", "") + "\n"
        return text.strip()
    except Exception as e:
        st.error(f"Error reading PPTX: {str(e)}")
        return ""

def extract_text_from_txt(file):
    """Extract text from TXT file"""
    try:
        return file.read().decode('utf-8').strip()
    except Exception as e:
        st.error(f"Error reading TXT: {str(e)}")
        return ""

# Add image Q&A support
from PIL import Image

def extract_image_bytes(file):
    """Return image bytes if file is an image, else None."""
    try:
        file.seek(0)
        img = Image.open(file)
        img.verify()  # Verify it's an image
        file.seek(0)
        return file.read()
    except Exception:
        return None

# Define a dictionary for text extractors
text_extractors: Dict[str, Callable[[io.BytesIO], str]] = {
    "application/pdf": extract_text_from_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_text_from_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_text_from_pptx,
    "text/plain": extract_text_from_txt,
}

# Updated process_document to handle None file_type
def process_document(uploaded_file: Optional[io.BytesIO]) -> str:
    """Process uploaded document and extract text with enhanced error handling."""
    if not uploaded_file:
        logging.warning("No file uploaded.")
        return ""

    # Infer file type from file extension
    file_type, _ = mimetypes.guess_type(uploaded_file.name)
    if not file_type:
        logging.error("Could not determine file type.")
        st.error("Could not determine file type.")
        return ""

    extractor = text_extractors.get(file_type)

    if extractor:
        try:
            return extractor(uploaded_file)
        except Exception as e:
            logging.error(f"Error extracting text from file: {e}")
            st.error(f"Error extracting text from file: {e}")
            return ""
    else:
        logging.error(f"Unsupported file type: {file_type}")
        st.error(f"Unsupported file type: {file_type}")
        return ""

def split_text_into_chunks(text: str, max_tokens: int) -> list[str]:
    """Split text into smaller chunks that fit within the token limit."""
    words = text.split()
    chunks = []
    current_chunk = []
    current_length = 0

    for word in words:
        if current_length + len(word) + 1 > max_tokens:  # +1 for space
            chunks.append(" ".join(current_chunk))
            current_chunk = []
            current_length = 0

        current_chunk.append(word)
        current_length += len(word) + 1

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks

# Updated ask_groq function with compound backoff and progress bar
def ask_groq(question: str, document_text: str, api_key: str) -> str:
    """Ask question using Groq API with document context and robust error handling."""
    try:
        # Split the document into chunks
        max_tokens_per_chunk = 5000  # Adjust based on the model's token limit
        chunks = split_text_into_chunks(document_text, max_tokens_per_chunk)

        aggregated_response = ""

        # Add a progress bar for chunk processing
        progress_bar = st.progress(0)

        for i, chunk in enumerate(chunks):
            logging.info(f"Processing chunk {i + 1}/{len(chunks)}")

            # Initialize Groq client with only the API key
            client = Groq(api_key=api_key)

            # Create a comprehensive prompt with document context
            prompt = f"""You are a helpful AI assistant. Answer the question based on the provided document context. If the answer cannot be found in the document, say so clearly.

DOCUMENT CONTENT:
{chunk}

QUESTION: {question}

Please provide a clear, accurate answer based on the document content above. If the information needed to answer the question is not in the document, please state that clearly."""

            retry_count = 0
            max_retries = 5
            backoff_factor = 2
            delay = 1  # Initial delay in seconds

            while retry_count < max_retries:
                try:
                    response = client.chat.completions.create(
                        messages=[
                            {
                                "role": "user",
                                "content": prompt
                            }
                        ],
                        model="llama3-8b-8192",
                        max_tokens=1000,
                        temperature=0.1,
                    )

                    # Safely append response content
                    if response.choices and response.choices[0].message.content:
                        aggregated_response += response.choices[0].message.content + "\n"

                    break  # Exit retry loop on success

                except Exception as e:
                    if "429 Too Many Requests" in str(e):
                        retry_count += 1
                        logging.info(f"Retrying request to /openai/v1/chat/completions in {delay} seconds (Attempt {retry_count}/{max_retries})")
                        time.sleep(delay)
                        delay *= backoff_factor  # Compound backoff
                    else:
                        logging.error(f"Error during Groq API call: {e}")
                        return f"Error: {str(e)}"

            if retry_count == max_retries:
                logging.error("Max retries reached. Unable to process chunk.")
                return "Error: Max retries reached. Unable to process chunk."

            # Update progress bar
            progress_bar.progress((i + 1) / len(chunks))

        return aggregated_response.strip()

    except Exception as e:
        logging.error(f"Error during Groq API call: {e}")
        return f"Error: {str(e)}"

def initialize_groq_client(api_key: str) -> Optional[Groq]:
    """Initialize the Groq client with error handling and detailed logging."""
    try:
        # Initialize the Groq client with only the API key
        client = Groq(api_key=api_key)
        logging.info("Groq client initialized successfully.")
        return client
    except TypeError as e:
        if "proxies" in str(e):
            logging.error("The Groq client received an unsupported argument 'proxies'. Please ensure the SDK is up-to-date.")
        else:
            logging.error(f"Failed to initialize Groq client: {e}")
        return None
    except Exception as e:
        logging.error(f"Unexpected error during Groq client initialization: {e}")
        return None

# Add support for Compound Beta and DeepSeek models
def ask_groq_with_model(question: str, document_text: str, api_key: str, model: str) -> str:
    """Ask question using a specific Groq model with document context."""
    try:
        max_tokens_per_chunk = 5000
        chunks = split_text_into_chunks(document_text, max_tokens_per_chunk)
        aggregated_response = ""
        progress_bar = st.progress(0)
        for i, chunk in enumerate(chunks):
            logging.info(f"Processing chunk {i + 1}/{len(chunks)} with model {model}")
            client = Groq(api_key=api_key)
            prompt = f"""You are a helpful AI assistant. Answer the question based on the provided document context. If the answer cannot be found in the document, say so clearly.\n\nDOCUMENT CONTENT:\n{chunk}\n\nQUESTION: {question}\n\nPlease provide a clear, accurate answer based on the document content above. If the information needed to answer the question is not in the document, please state that clearly."""
            retry_count = 0
            max_retries = 10 if "compound-beta" in model else 5
            backoff_factor = 2
            delay = 2 if "compound-beta" in model else 1
            while retry_count < max_retries:
                try:
                    response = client.chat.completions.create(
                        messages=[{"role": "user", "content": prompt}],
                        model=model,
                        max_tokens=1000,
                        temperature=0.1,
                    )
                    if response.choices and response.choices[0].message.content:
                        aggregated_response += response.choices[0].message.content + "\n"
                    break
                except Exception as e:
                    if "429 Too Many Requests" in str(e):
                        retry_count += 1
                        logging.info(f"Rate limit reached for model {model}. Retrying in {delay} seconds (Attempt {retry_count}/{max_retries})")
                        st.warning(f"Rate limit reached for model {model}. Retrying in {delay} seconds (Attempt {retry_count}/{max_retries})")
                        time.sleep(delay)
                        delay *= backoff_factor
                    else:
                        logging.error(f"Error during Groq API call: {e}")
                        return f"Error: {str(e)}"
            if retry_count == max_retries:
                logging.error("Max retries reached. Unable to process chunk.")
                return "Error: Max retries reached. Unable to process chunk."
            progress_bar.progress((i + 1) / len(chunks))
        return aggregated_response.strip()
    except Exception as e:
        logging.error(f"Error during Groq API call: {e}")
        return f"Error: {str(e)}"

# Move ask_groq_vision definition above the main UI logic so it is always defined before use.
def ask_groq_vision(question: str, image_bytes: bytes, api_key: str, model: str = "meta-llama/llama-4-scout-17b-16e-instruct") -> str:
    """Call Groq Vision API for image Q&A using base64-encoded image in message content. Model can be 'scout' or 'maverick'."""
    import base64
    from groq import Groq
    # Encode image as base64
    base64_image = base64.b64encode(image_bytes).decode('utf-8')
    client = Groq(api_key=api_key)
    try:
        chat_completion = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": question},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}",
                            },
                        },
                    ],
                }
            ],
            model=model,
        )
        return chat_completion.choices[0].message.content or "[No answer returned]"
    except Exception as e:
        return f"Groq Vision API error: {e}"

# --- Universal Q&A handler ---
def handle_upload_and_qa(uploaded_file, question, api_key):
    """Automatically detect file type and use the correct model/API for Q&A."""
    if not uploaded_file or not question or not api_key:
        return None, None, None
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    # Image types
    if file_extension in [".jpg", ".jpeg", ".png", ".bmp", ".gif", ".webp"]:
        image_bytes = extract_image_bytes(uploaded_file)
        if image_bytes:
            vision_model = "meta-llama/llama-4-scout-17b-16e-instruct"
            answer = ask_groq_vision(question, image_bytes, api_key, model=vision_model)
            return answer, vision_model, "image"
        else:
            return "Failed to read image file.", None, None
    # Document types
    elif file_extension in [".pdf", ".docx", ".pptx", ".txt"]:
        document_text = process_document(uploaded_file)
        if document_text:
            answer = ask_groq(question, document_text, api_key)
            return answer, "llama3-8b-8192", "text"
        else:
            return "Failed to extract text from document.", None, None
    # Fallback: try as image, then as text
    image_bytes = extract_image_bytes(uploaded_file)
    if image_bytes:
        vision_model = "meta-llama/llama-4-scout-17b-16e-instruct"
        answer = ask_groq_vision(question, image_bytes, api_key, model=vision_model)
        return answer, vision_model, "image"
    document_text = process_document(uploaded_file)
    if document_text:
        answer = ask_groq(question, document_text, api_key)
        return answer, "llama3-8b-8192", "text"
    return "Unsupported file type or failed to process.", None, None

# Initialize variables to avoid unbound errors
summary = "No summary available."
key_points = "No key points available."
main_topic = "No main topic available."
document_type = "Unknown"

# Main UI
st.title("🤖 RAQ - Pure Groq API")
st.markdown("Simple document Q&A powered entirely by Groq API")

# Sidebar
st.sidebar.header("🔑 Configuration")

# API Key input
api_key = st.sidebar.text_input("Enter your Groq API Key:", type="password", help="Get your free API key from https://console.groq.com")

if api_key:
    st.session_state.groq_client = initialize_groq_client(api_key)
    if st.session_state.groq_client:
        st.sidebar.success("✅ API Key Set!")
    else:
        st.sidebar.error("❌ Failed to initialize Groq client. Check logs for details.")

# Sidebar
st.sidebar.header("📄 Document Upload")
uploaded_file = st.sidebar.file_uploader(
    "Choose a file", 
    type=['pdf', 'docx', 'pptx', 'txt', 'jpg', 'jpeg', 'png'],
    help="Upload PDF, DOCX, PPTX, TXT, or image files"
)

# Document type detection
if uploaded_file:
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    document_type = "Unknown"
    if file_extension in [".pdf"]:
        document_type = "PDF Document"
    elif file_extension in [".docx"]:
        document_type = "Word Document"
    elif file_extension in [".pptx"]:
        document_type = "PowerPoint Presentation"
    elif file_extension in [".txt"]:
        document_type = "Text File"
    elif file_extension in [".jpg", ".jpeg", ".png"]:
        document_type = "Image"
    st.sidebar.info(f"📄 Detected Document Type: {document_type}")

# Check file size limit
if uploaded_file:
    if uploaded_file.size > 5 * 1024 * 1024:  # 5 MB limit
        st.sidebar.error("❌ File size exceeds the 5 MB limit. Please upload a smaller file.")
        logging.warning("Uploaded file exceeds size limit.")
    else:
        with st.spinner("📖 Processing document..."):
            document_text = process_document(uploaded_file)
            if document_text:
                st.session_state.document_text = document_text
                st.sidebar.success(f"✅ Document processed! ({len(document_text)} characters)")
                logging.info(f"Document processed successfully. Length: {len(document_text)} characters.")

                # Perform summary analysis
                with st.spinner("🤔 Generating summary analysis..."):
                    summary_prompt = "Please provide a comprehensive summary of this document."
                    key_points_prompt = "What are the main key points or takeaways from this document?"
                    main_topic_prompt = "What is the main topic or subject of this document?"

                    summary = ask_groq(summary_prompt, document_text, api_key)
                    key_points = ask_groq(key_points_prompt, document_text, api_key)
                    main_topic = ask_groq(main_topic_prompt, document_text, api_key)

                    # Display suggested prompts
                    st.sidebar.markdown("### Suggested Prompts")
                    st.sidebar.button("📝 Summary", on_click=lambda: st.session_state.update({'question': summary_prompt}))
                    st.sidebar.button("🔍 Key Points", on_click=lambda: st.session_state.update({'question': key_points_prompt}))
                    st.sidebar.button("❓ Main Topic", on_click=lambda: st.session_state.update({'question': main_topic_prompt}))

            else:
                st.sidebar.error("❌ Failed to extract text from document")
                logging.error("Failed to extract text from document.")

# User feedback mechanism
if 'last_answer' in st.session_state:
    st.markdown("### 💬 Feedback on Answer")
    feedback = st.radio("Was this answer helpful?", ["Yes", "No"], key="feedback_radio")
    if st.button("Submit Feedback"):
        st.session_state.feedback.append({
            "question": st.session_state.get("question", ""),
            "answer": st.session_state.last_answer,
            "feedback": feedback
        })
        st.success("Thank you for your feedback!")

# Updated export results logic
if st.button("📥 Export Analysis Results"):
    analysis_results = f"Summary: {summary}\nKey Points: {key_points}\nMain Topic: {main_topic}"
    st.download_button(
        label="Download Results",
        data=analysis_results,
        file_name="analysis_results.txt",
        mime="text/plain"
    )

# Updated session history logic
if st.button("📜 View Session History"):
    if st.session_state.history:
        history_df = pd.DataFrame(st.session_state.history)
        st.write(history_df)
    else:
        st.info("No session history available.")

# Updated advanced analytics logic
if st.session_state.get("document_text"):
    document_text = st.session_state.document_text
    st.sidebar.markdown("### 📊 Advanced Analytics")
    sentiment = "Positive" if "good" in document_text.lower() else "Neutral"
    keywords = ", ".join(set(document_text.split()[:5]))  # Simplified keyword extraction
    st.sidebar.metric("Sentiment", sentiment)
    st.sidebar.metric("Keywords", keywords)

# Performance metrics
st.sidebar.markdown("### ⏱️ Performance Metrics")
processing_time = time.time() - st.session_state.get("start_time", time.time())
st.sidebar.metric("Processing Time", f"{processing_time:.2f} seconds")

# Floating action button
st.markdown(
    """<style>
    .floating-button {
        position: fixed;
        bottom: 20px;
        right: 20px;
        background-color: #4CAF50;
        color: white;
        border: none;
        border-radius: 50%;
        width: 60px;
        height: 60px;
        font-size: 24px;
        text-align: center;
        line-height: 60px;
        box-shadow: 0px 4px 6px rgba(0, 0, 0, 0.1);
        cursor: pointer;
    }
    </style>
    <button class="floating-button">+</button>
    """,
    unsafe_allow_html=True
)

# Main content area
col1, col2 = st.columns([2, 1])

with col1:
    st.header("💬 Ask Questions")
    if not api_key:
        st.info("👈 Please enter your Groq API key in the sidebar to get started")
    elif not uploaded_file:
        st.info("👈 Please upload a document or image in the sidebar first")
    else:
        question = st.text_input("What would you like to know about the uploaded file?")
        if question:
            with st.spinner("🤔 Processing..."):
                answer, used_model, file_type = handle_upload_and_qa(uploaded_file, question, api_key)
            if file_type == "image":
                st.markdown(f"### 👁️ Vision Q&A Answer ({used_model}):")
                st.markdown(answer)
                image_bytes = extract_image_bytes(uploaded_file)
                if image_bytes:
                    st.image(image_bytes, caption="Uploaded Image", use_container_width=True)
            elif file_type == "text":
                st.markdown(f"### 💡 Text Q&A Answer ({used_model}):")
                st.markdown(answer)
            else:
                st.error(answer)

    # Move analysis results to the main body of the app
    st.markdown("### 📊 Analysis Results")
    st.markdown(f"**Summary:** {summary}")
    st.markdown(f"**Key Points:** {key_points}")
    st.markdown(f"**Main Topic:** {main_topic}")

with col2:
    st.header("ℹ️ About")
    st.markdown("""
    **RAQ** uses only Groq API calls for:
    
    ✅ **Simple & Fast**  
    No complex embeddings or vector databases
    
    ✅ **Reliable**  
    No dependency conflicts
    
    ✅ **Pure Groq**  
    Leverages Groq's fast inference
    
    ✅ **Multi-format**  
    PDF, DOCX, PPTX, TXT support
    """)
    
    if st.session_state.document_text:
        st.markdown("### 📊 Document Stats")
        st.metric("Characters", len(st.session_state.document_text))
        st.metric("Words", len(st.session_state.document_text.split()))

# Footer
st.markdown("---")
st.markdown("**Powered by Groq API** • [Get API Key](https://console.groq.com) • Simple & Fast RAG")

# Add a new tab for advanced reasoning chatbot
st.sidebar.header("🧠 Advanced Reasoning Chatbot")
selected_model = st.sidebar.selectbox(
    "Choose a model for advanced reasoning:",
    ["compound-beta", "compound-beta-mini", "deepseek-r1-distill-llama-70b"],
    help="Select a model for advanced reasoning tasks."
)

if selected_model:
    st.header("🧠 Advanced Reasoning Chatbot")
    question = st.text_input(
        "Ask a question for advanced reasoning:",
        placeholder="e.g., What are the implications of quantum computing?"
    )

    if question and api_key:
        with st.spinner("🤔 Thinking..."):
            answer = ask_groq_with_model(question, st.session_state.document_text, api_key, selected_model)
        st.markdown("### 💡 Answer:")
        st.markdown(answer)

# Refactor Vision Q&A capabilities into a dedicated function
def vision_qa_section(api_key):
    """Encapsulates Vision Q&A logic for image uploads and question answering."""
    st.header("👁️ Vision Q&A")

    # File uploader for Vision Q&A
    vision_uploaded_file = st.file_uploader(
        "Upload an image for Vision Q&A", 
        type=['jpg', 'jpeg', 'png', 'bmp', 'gif', 'webp'],
        help="Upload image files for Vision Q&A"
    )

    if vision_uploaded_file:
        with st.spinner("📖 Processing image..."):
            image_bytes = extract_image_bytes(vision_uploaded_file)
            if image_bytes:
                st.success("✅ Image uploaded successfully!")
                question = st.text_input("What would you like to know about the image?")
                if question:
                    with st.spinner("🤔 Processing Vision Q&A..."):
                        vision_model = "meta-llama/llama-4-scout-17b-16e-instruct"
                        answer = ask_groq_vision(question, image_bytes, api_key, model=vision_model)
                        st.markdown(f"### 👁️ Vision Q&A Answer ({vision_model}):")
                        st.markdown(answer)
                        st.image(image_bytes, caption="Uploaded Image", use_container_width=True)
            else:
                st.error("❌ Failed to read image file.")

# Call the refactored Vision Q&A section in the main layout
vision_qa_section(api_key)
